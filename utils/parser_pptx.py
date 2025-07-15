# utils/parser_pptx.py
from pptx import Presentation

def parse_pptx_status(file_path):
    prs = Presentation(file_path)
    text_runs = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text_runs.append(shape.text)
    return {
        "parsed": {},
        "raw_text": "\n".join(text_runs)
    }